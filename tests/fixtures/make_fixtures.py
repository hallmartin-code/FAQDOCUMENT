"""Regenerate the binary deck fixtures.

Run from the repo root:  python tests/fixtures/make_fixtures.py

The generated files are committed so the test suite never depends on this script, but
keeping the generator around documents exactly what the fixtures contain.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import landscape, letter
from reportlab.pdfgen import canvas

HERE = Path(__file__).parent

SLIDES: list[tuple[str, list[str]]] = [
    (
        "Helion Bio",
        [
            "Reprogramming macrophages to clear solid tumors",
            "Seed round | $4M | San Diego, CA",
        ],
    ),
    (
        "The Problem",
        [
            "Checkpoint inhibitors fail in 80% of solid tumor patients.",
            "Tumor-associated macrophages actively suppress T-cell response.",
            "No approved therapy targets the macrophage compartment directly.",
        ],
    ),
    (
        "Our Solution",
        [
            "HLN-101 is a small molecule SIRP-alpha antagonist.",
            "Repolarizes M2 macrophages to an M1 phenotype in vitro.",
            "Oral dosing; no cold chain required.",
        ],
    ),
    (
        "Team",
        [
            "Dr. Marisol Reyes, CEO - 12 years at Genentech, led two INDs.",
            "Dr. Peter Kwan, CSO - Professor of Immunology, UCSD.",
            "Commercial lead role is currently open.",
        ],
    ),
    (
        "The Ask",
        [
            "Raising $4M seed to complete IND-enabling toxicology.",
            "18 month runway to a Phase 1 start.",
            "No lead investor committed to date.",
        ],
    ),
]


def build_pdf(target: Path) -> None:
    """Write a five-page landscape PDF deck."""
    page = landscape(letter)
    pdf = canvas.Canvas(str(target), pagesize=page)
    width, height = page
    for title, bullets in SLIDES:
        pdf.setFont("Helvetica-Bold", 30)
        pdf.drawString(60, height - 100, title)
        pdf.setFont("Helvetica", 16)
        y = height - 170
        for bullet in bullets:
            pdf.drawString(70, y, bullet)
            y -= 34
        pdf.showPage()
    pdf.save()


def build_pptx(target: Path) -> None:
    """Write the same five slides as a PPTX, with speaker notes on every slide."""
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    for number, (title, bullets) in enumerate(SLIDES, start=1):
        slide = presentation.slides.add_slide(blank)
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8.5), Inches(1.0))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].runs[0].font.size = Pt(32)

        body = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(8.5), Inches(4.0))
        body_frame = body.text_frame
        body_frame.text = bullets[0]
        for bullet in bullets[1:]:
            body_frame.add_paragraph().text = bullet

        slide.notes_slide.notes_text_frame.text = (
            f"Speaker notes for slide {number}: emphasize the {title.lower()} narrative."
        )
    presentation.save(str(target))


def build_image_heavy_pdf(target: Path) -> None:
    """Write a three-page deck that says almost nothing in text.

    Page 1 is a bar chart with a two-character axis label, page 2 a scatter of circles
    with no text at all, page 3 a dense grid. Every page falls under the 20-character
    image-dominant threshold and over the vector-path count that reads as a chart, which
    is exactly the deck the ingest layer has to notice and send as pictures.
    """
    page = landscape(letter)
    pdf = canvas.Canvas(str(target), pagesize=page)
    width, height = page

    # Page 1 — bar chart.
    for i in range(24):
        bar_height = 40 + (i * 37) % 260
        pdf.rect(80 + i * 26, 120, 18, bar_height, stroke=1, fill=1)
    pdf.setFont("Helvetica", 9)
    pdf.drawString(80, 100, "$M")
    pdf.showPage()

    # Page 2 — scatter, no text whatsoever.
    for i in range(60):
        pdf.circle(120 + (i * 53) % 700, 120 + (i * 91) % 380, 6, stroke=1, fill=0)
    pdf.showPage()

    # Page 3 — grid.
    for i in range(20):
        pdf.line(60, 100 + i * 22, width - 60, 100 + i * 22)
    for i in range(20):
        pdf.line(60 + i * 34, 100, 60 + i * 34, height - 100)
    pdf.showPage()

    pdf.save()

if __name__ == "__main__":
    build_pdf(HERE / "sample_deck.pdf")
    build_pptx(HERE / "sample_deck.pptx")
    build_image_heavy_pdf(HERE / "image_heavy_deck.pdf")
    for name in ("sample_deck.pdf", "sample_deck.pptx", "image_heavy_deck.pdf"):
        print(f"wrote {HERE / name}")
