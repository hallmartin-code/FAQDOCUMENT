"""PPTX ingestion.

Text and speaker notes come from python-pptx. Visuals require a LibreOffice conversion to
PDF followed by PyMuPDF rasterization; when `soffice` is absent we degrade to text-only and
say so loudly, because chart-heavy slides will be under-read.
"""

from __future__ import annotations

import base64
import shutil
import subprocess
import tempfile
from pathlib import Path

import pymupdf  # PyMuPDF; the legacy `fitz` alias is deprecated as of 1.28
from pptx import Presentation
from pptx.shapes.base import BaseShape

from deckpager.errors import IngestError
from deckpager.ingest.models import Deck, Slide, SlideAsset, normalize_text
from deckpager.ingest.pdf import first_line_title, render_page_png

#: Seconds to wait for LibreOffice before giving up and degrading to text-only.
SOFFICE_TIMEOUT_S = 180


def _shape_text(shape: BaseShape) -> list[str]:
    """Pull text out of a shape, descending into groups and tables."""
    parts: list[str] = []
    if shape.shape_type is not None and shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for member in shape.shapes:  # type: ignore[attr-defined]
            parts.extend(_shape_text(member))
        return parts
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:  # type: ignore[attr-defined]
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
        return parts
    if shape.has_text_frame and shape.text_frame.text.strip():  # type: ignore[attr-defined]
        parts.append(shape.text_frame.text)  # type: ignore[attr-defined]
    return parts


def find_soffice() -> str | None:
    """Locate the LibreOffice CLI, checking PATH then the usual Windows install roots."""
    for name in ("soffice", "soffice.exe"):
        found = shutil.which(name)
        if found:
            return found
    for candidate in (
        Path(r"C:\Program Files\LibreOffice\program\soffice.exe"),
        Path(r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"),
    ):
        if candidate.is_file():
            return str(candidate)
    return None


def _rasterize_via_libreoffice(path: Path, soffice: str) -> list[bytes]:
    """Convert the PPTX to PDF with LibreOffice, then rasterize each page to PNG."""
    with tempfile.TemporaryDirectory(prefix="deckpager-pptx-") as tmp:
        tmp_dir = Path(tmp)
        result = subprocess.run(  # noqa: S603 - soffice path resolved above, args are fixed
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(tmp_dir), str(path)],
            capture_output=True,
            timeout=SOFFICE_TIMEOUT_S,
            check=False,
        )
        produced = sorted(tmp_dir.glob("*.pdf"))
        if result.returncode != 0 or not produced:
            detail = result.stderr.decode("utf-8", "replace").strip() or "no PDF produced"
            raise RuntimeError(f"LibreOffice conversion failed: {detail}")
        with pymupdf.open(produced[0]) as document:
            return [render_page_png(page) for page in document]


def load_pptx(path: Path, *, want_images: bool) -> Deck:
    """Read a PPTX into a Deck, rasterizing slides when LibreOffice is available."""
    try:
        presentation = Presentation(str(path))
    except Exception as exc:  # python-pptx raises package-specific errors for bad files
        raise IngestError(f"{path.name} is not a readable PPTX: {exc}") from exc

    warnings: list[str] = []
    slides: list[Slide] = []
    for number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            parts.extend(_shape_text(shape))
        body = normalize_text("\n".join(parts))

        # Prefer the title placeholder; many real decks use plain text boxes on blank
        # layouts, so fall back to the first line the same way the PDF path does.
        title_shape = slide.shapes.title
        title = title_shape.text_frame.text.strip()[:200] if title_shape is not None else ""
        if not title:
            title = first_line_title(body) or ""

        notes: str | None = None
        if slide.has_notes_slide:
            raw_notes = slide.notes_slide.notes_text_frame.text
            notes = normalize_text(raw_notes) or None

        slides.append(
            Slide(
                index=number,
                title=title or None,
                text=body,
                notes=notes,
                asset=None,
            )
        )

    if not slides:
        raise IngestError(f"{path.name} contains no slides.")

    if want_images:
        soffice = find_soffice()
        if soffice is None:
            warnings.append(
                "LibreOffice (soffice) not found on PATH — slides were not rasterized. "
                "Chart-heavy and image-only slides may be under-read; install LibreOffice "
                "or pass --no-images to silence this."
            )
        else:
            try:
                pages = _rasterize_via_libreoffice(path, soffice)
            except (RuntimeError, subprocess.TimeoutExpired, OSError) as exc:
                warnings.append(
                    f"Slide rasterization failed ({exc}); analysis proceeds on text only "
                    f"and chart-heavy slides may be under-read."
                )
            else:
                if len(pages) != len(slides):
                    warnings.append(
                        f"LibreOffice produced {len(pages)} pages for {len(slides)} slides; "
                        f"images were matched positionally and may be offset."
                    )
                for slide_model, png in zip(slides, pages, strict=False):
                    slide_model.asset = SlideAsset(
                        data_b64=base64.standard_b64encode(png).decode("ascii")
                    )

    return Deck(
        source_path=path,
        source_format="pptx",
        slides=slides,
        raw_pdf_b64=None,
        warnings=warnings,
    )
